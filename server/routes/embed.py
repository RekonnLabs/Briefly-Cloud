"""
Document embedding routes for Briefly Cloud
Handles document indexing and vector embedding
"""

from fastapi import APIRouter, HTTPException, BackgroundTasks, Depends
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import os
import logging
import asyncio
from datetime import datetime
from dotenv import load_dotenv
from supabase import create_client, Client
import httpx
# Google API imports (optional for serverless deployment)
try:
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    GOOGLE_API_AVAILABLE = True
except ImportError as e:
    print(f"Google API libraries not available: {e}")
    GOOGLE_API_AVAILABLE = False
    # Create dummy classes for compatibility
    class Credentials:
        def __init__(self, *args, **kwargs):
            pass
    def build(*args, **kwargs):
        raise HTTPException(status_code=501, detail="Google API not available in this deployment")
import io
import tempfile

# Document processing imports
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
import json

# ML imports (optional for serverless deployment)
try:
    from sentence_transformers import SentenceTransformer
    import chromadb
    from chromadb.config import Settings
    ML_LIBRARIES_AVAILABLE = True
except ImportError as e:
    print(f"ML libraries not available: {e}")
    ML_LIBRARIES_AVAILABLE = False
    # Create dummy classes for compatibility
    class SentenceTransformer:
        def __init__(self, *args, **kwargs):
            pass
        def encode(self, *args, **kwargs):
            raise HTTPException(status_code=501, detail="ML libraries not available in this deployment")
    
    class chromadb:
        @staticmethod
        def Client(*args, **kwargs):
            raise HTTPException(status_code=501, detail="ChromaDB not available in this deployment")
    
    class Settings:
        pass

# Import usage limits
from utils.usage_limits import (
    check_and_increment_usage, 
    UsageLimitError,
    format_storage_size
)

# Load environment variables
load_dotenv()

# Configure logging
logger = logging.getLogger(__name__)

# Initialize Supabase client
# Import shared Supabase client
from utils.supabase_client import get_supabase_client

# Initialize embedding model
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# Initialize ChromaDB client for Chroma Cloud
chroma_client = None
try:
    if os.getenv("CHROMA_API_KEY") and os.getenv("CHROMA_TENANT_ID"):
        chroma_client = chromadb.CloudClient(
            api_key=os.getenv("CHROMA_API_KEY"),
            tenant=os.getenv("CHROMA_TENANT_ID"),
            database=os.getenv("CHROMA_DB_NAME", "Briefly Cloud")
        )
        logger.info("✅ ChromaDB Cloud client initialized successfully")
    else:
        logger.info("ℹ️ ChromaDB credentials not configured, embedding disabled")
except Exception as e:
    logger.info(f"ℹ️ ChromaDB unavailable, embedding disabled")
    chroma_client = None

router = APIRouter(prefix="/embed", tags=["embedding"])

# Data models
class EmbedRequest(BaseModel):
    user_id: str
    source: str  # 'google' or 'microsoft'
    file_ids: Optional[List[str]] = None  # Specific files to process

class EmbedStatus(BaseModel):
    status: str  # 'pending', 'processing', 'completed', 'failed'
    progress: float  # 0.0 to 1.0
    message: str
    files_processed: int
    total_files: int

class DocumentChunk(BaseModel):
    content: str
    metadata: Dict[str, Any]
    chunk_index: int

# Supported file types and their MIME types
SUPPORTED_MIME_TYPES = {
    'application/pdf': 'pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'text/plain': 'txt',
    'application/json': 'json'
}

@router.post("/start")
async def start_embedding(
    request: EmbedRequest,
    background_tasks: BackgroundTasks
):
    """Start document embedding process"""
    try:
        # Enforce document upload usage limits
        usage_data = await check_and_increment_usage(
            supabase,
            request.user_id,
            'documents',
            'document_upload',
            increment=len(request.file_ids),
            event_data={
                'source': request.source,
                'file_count': len(request.file_ids),
                'file_ids': request.file_ids
            }
        )
        
        # Validate user and storage connection
        token_data = get_supabase_client().table('oauth_tokens').select('*').eq('user_id', request.user_id).eq('provider', request.source).execute()
        
        if not token_data.data:
            raise HTTPException(status_code=404, detail=f"{request.source.title()} storage not connected")
        
        # Create job log entry
        job_data = {
            "user_id": request.user_id,
            "job_type": "document_embedding",
            "status": "pending",
            "input_data": request.dict()
        }
        
        job_result = get_supabase_client().table('job_logs').insert(job_data).execute()
        job_id = job_result.data[0]['id']
        
        # Start background processing
        background_tasks.add_task(
            process_documents_background,
            request.user_id,
            request.source,
            job_id,
            request.file_ids
        )
        
        return {
            "message": "Document embedding started",
            "job_id": job_id,
            "status": "pending"
        }
        
    except Exception as e:
        logger.error(f"Embedding start error: {e}")
        raise HTTPException(status_code=500, detail="Failed to start embedding")

@router.get("/status/{job_id}")
async def get_embedding_status(job_id: str):
    """Get embedding job status"""
    try:
        job_data = get_supabase_client().table('job_logs').select('*').eq('id', job_id).execute()
        
        if not job_data.data:
            raise HTTPException(status_code=404, detail="Job not found")
            
        job = job_data.data[0]
        output_data = job.get('output_data', {})
        
        return EmbedStatus(
            status=job['status'],
            progress=output_data.get('progress', 0.0),
            message=output_data.get('message', ''),
            files_processed=output_data.get('files_processed', 0),
            total_files=output_data.get('total_files', 0)
        )
        
    except Exception as e:
        logger.error(f"Status check error: {e}")
        raise HTTPException(status_code=500, detail="Failed to get status")

async def process_documents_background(
    user_id: str,
    source: str,
    job_id: str,
    file_ids: Optional[List[str]] = None
):
    """Background task to process and embed documents"""
    try:
        # Update job status
        await update_job_status(job_id, "processing", {"message": "Starting document processing..."})
        
        # Get OAuth tokens
        token_data = get_supabase_client().table('oauth_tokens').select('*').eq('user_id', user_id).eq('provider', source).execute()
        token_info = token_data.data[0]
        
        # Get files from cloud storage
        if source == 'google':
            files = await get_google_drive_files(token_info, file_ids)
        else:
            files = await get_onedrive_files(token_info, file_ids)
        
        # Filter and validate supported files
        supported_files = []
        skipped_files = []
        
        for file_info in files:
            mime_type = file_info.get('mimeType', '')
            if mime_type in SUPPORTED_MIME_TYPES:
                supported_files.append(file_info)
            else:
                skipped_files.append({
                    'name': file_info['name'],
                    'mime_type': mime_type,
                    'reason': 'Unsupported file type'
                })
        
        total_files = len(supported_files)
        await update_job_status(job_id, "processing", {
            "message": f"Found {total_files} supported files to process ({len(skipped_files)} skipped)",
            "total_files": total_files,
            "files_processed": 0,
            "progress": 0.0,
            "skipped_files": skipped_files
        })
        
        # Get or create user's collection in ChromaDB
        collection_name = f"user_{user_id}"
        try:
            collection = chroma_client.get_collection(collection_name)
        except:
            collection = chroma_client.create_collection(collection_name)
        
        # Process each supported file
        processed_count = 0
        for i, file_info in enumerate(supported_files):
            try:
                await update_job_status(job_id, "processing", {
                    "message": f"Processing {file_info['name']}...",
                    "files_processed": processed_count,
                    "total_files": total_files,
                    "progress": i / total_files
                })
                
                # Download and process file
                if source == 'google':
                    content = await download_google_file(token_info, file_info)
                else:
                    content = await download_onedrive_file(token_info, file_info)
                
                if content:
                    # Extract text and create chunks
                    chunks = await extract_and_chunk_text(content, file_info)
                    
                    # Generate embeddings and store
                    await store_document_chunks(user_id, file_info, chunks, collection)
                    
                    # Store file metadata
                    await store_file_metadata(user_id, file_info, source, len(chunks))
                    
                    processed_count += 1
                
            except Exception as e:
                logger.error(f"Error processing file {file_info['name']}: {e}")
                continue
        
        # Complete job
        await update_job_status(job_id, "completed", {
            "message": f"Successfully processed {processed_count} files",
            "files_processed": processed_count,
            "total_files": total_files,
            "progress": 1.0
        })
        
    except Exception as e:
        logger.error(f"Background processing error: {e}")
        await update_job_status(job_id, "failed", {
            "message": f"Processing failed: {str(e)}",
            "progress": 0.0
        })

async def update_job_status(job_id: str, status: str, output_data: Dict[str, Any]):
    """Update job status in database"""
    try:
        update_data = {
            "status": status,
            "output_data": output_data
        }
        
        if status == "completed" or status == "failed":
            update_data["completed_at"] = datetime.utcnow().isoformat()
            
        get_supabase_client().table('job_logs').update(update_data).eq('id', job_id).execute()
        
    except Exception as e:
        logger.error(f"Failed to update job status: {e}")

async def get_google_drive_files(token_info: Dict, file_ids: Optional[List[str]] = None):
    """Get files from Google Drive"""
    try:
        credentials = Credentials(
            token=token_info['access_token'],
            refresh_token=token_info['refresh_token']
        )
        
        service = build('drive', 'v3', credentials=credentials)
        
        if file_ids:
            # Get specific files
            files = []
            for file_id in file_ids:
                file_info = service.files().get(
                    fileId=file_id,
                    fields="id, name, mimeType, size, modifiedTime"
                ).execute()
                files.append(file_info)
        else:
            # Get all supported files
            query = " or ".join([f"mimeType='{mime}'" for mime in SUPPORTED_MIME_TYPES.keys()])
            
            results = service.files().list(
                q=query,
                pageSize=100,
                fields="files(id, name, mimeType, size, modifiedTime)"
            ).execute()
            
            files = results.get('files', [])
        
        return files
        
    except Exception as e:
        logger.error(f"Google Drive files error: {e}")
        return []

async def get_onedrive_files(token_info: Dict, file_ids: Optional[List[str]] = None):
    """Get files from OneDrive"""
    try:
        headers = {"Authorization": f"Bearer {token_info['access_token']}"}
        
        if file_ids:
            # Get specific files
            files = []
            async with httpx.AsyncClient() as client:
                for file_id in file_ids:
                    response = await client.get(
                        f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}",
                        headers=headers
                    )
                    if response.status_code == 200:
                        files.append(response.json())
        else:
            # Get all files
            async with httpx.AsyncClient() as client:
                response = await client.get(
                    "https://graph.microsoft.com/v1.0/me/drive/root/children?$filter=file ne null&$top=100",
                    headers=headers
                )
                
                if response.status_code == 200:
                    data = response.json()
                    files = data.get('value', [])
                else:
                    files = []
        
        return files
        
    except Exception as e:
        logger.error(f"OneDrive files error: {e}")
        return []

async def download_google_file(token_info: Dict, file_info: Dict) -> Optional[bytes]:
    """Download file content from Google Drive"""
    try:
        credentials = Credentials(
            token=token_info['access_token'],
            refresh_token=token_info['refresh_token']
        )
        
        service = build('drive', 'v3', credentials=credentials)
        
        # Download file
        file_content = service.files().get_media(fileId=file_info['id']).execute()
        return file_content
        
    except Exception as e:
        logger.error(f"Google file download error: {e}")
        return None

async def download_onedrive_file(token_info: Dict, file_info: Dict) -> Optional[bytes]:
    """Download file content from OneDrive"""
    try:
        headers = {"Authorization": f"Bearer {token_info['access_token']}"}
        
        async with httpx.AsyncClient() as client:
            response = await client.get(
                f"https://graph.microsoft.com/v1.0/me/drive/items/{file_info['id']}/content",
                headers=headers
            )
            
            if response.status_code == 200:
                return response.content
            else:
                return None
                
    except Exception as e:
        logger.error(f"OneDrive file download error: {e}")
        return None

async def extract_and_chunk_text(content: bytes, file_info: Dict) -> List[DocumentChunk]:
    """Extract text from file content and create chunks"""
    try:
        mime_type = file_info.get('mimeType', '')
        file_type = SUPPORTED_MIME_TYPES.get(mime_type, 'unknown')
        
        # Extract text based on file type
        if file_type == 'pdf':
            text = extract_pdf_text(content)
        elif file_type == 'docx':
            text = extract_docx_text(content)
        elif file_type == 'pptx':
            text = extract_pptx_text(content)
        elif file_type == 'xlsx':
            text = extract_xlsx_text(content)
        elif file_type == 'txt':
            text = content.decode('utf-8')
        elif file_type == 'json':
            text = content.decode('utf-8')
        else:
            logger.warning(f"Unsupported file type: {mime_type}")
            return []
        
        # Create chunks (split by paragraphs, max 1000 chars)
        chunks = []
        paragraphs = text.split('\n\n')
        current_chunk = ""
        chunk_index = 0
        
        for paragraph in paragraphs:
            if len(current_chunk) + len(paragraph) > 1000 and current_chunk:
                # Save current chunk
                chunks.append(DocumentChunk(
                    content=current_chunk.strip(),
                    metadata={
                        "file_name": file_info['name'],
                        "file_id": file_info['id'],
                        "mime_type": mime_type,
                        "chunk_size": len(current_chunk)
                    },
                    chunk_index=chunk_index
                ))
                chunk_index += 1
                current_chunk = paragraph
            else:
                current_chunk += "\n\n" + paragraph if current_chunk else paragraph
        
        # Add final chunk
        if current_chunk.strip():
            chunks.append(DocumentChunk(
                content=current_chunk.strip(),
                metadata={
                    "file_name": file_info['name'],
                    "file_id": file_info['id'],
                    "mime_type": mime_type,
                    "chunk_size": len(current_chunk)
                },
                chunk_index=chunk_index
            ))
        
        return chunks
        
    except Exception as e:
        logger.error(f"Text extraction error: {e}")
        return []

def extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF content"""
    try:
        with tempfile.NamedTemporaryFile() as tmp_file:
            tmp_file.write(content)
            tmp_file.flush()
            
            text = ""
            with pdfplumber.open(tmp_file.name) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
            
            return text
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return ""

def extract_docx_text(content: bytes) -> str:
    """Extract text from DOCX content"""
    try:
        with tempfile.NamedTemporaryFile() as tmp_file:
            tmp_file.write(content)
            tmp_file.flush()
            
            doc = Document(tmp_file.name)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n\n"
            
            return text
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return ""

def extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX content"""
    try:
        with tempfile.NamedTemporaryFile() as tmp_file:
            tmp_file.write(content)
            tmp_file.flush()
            
            prs = Presentation(tmp_file.name)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n\n"
            
            return text
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return ""

def extract_xlsx_text(content: bytes) -> str:
    """Extract text from XLSX content"""
    try:
        with tempfile.NamedTemporaryFile() as tmp_file:
            tmp_file.write(content)
            tmp_file.flush()
            
            wb = openpyxl.load_workbook(tmp_file.name)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + "\n"
            
            return text
    except Exception as e:
        logger.error(f"XLSX extraction error: {e}")
        return ""

async def store_document_chunks(
    user_id: str,
    file_info: Dict,
    chunks: List[DocumentChunk],
    collection
):
    """Store document chunks with embeddings"""
    try:
        # Generate embeddings for all chunks
        chunk_texts = [chunk.content for chunk in chunks]
        embeddings = embedding_model.encode(chunk_texts).tolist()
        
        # Prepare data for ChromaDB
        ids = [f"{file_info['id']}_{chunk.chunk_index}" for chunk in chunks]
        metadatas = [chunk.metadata for chunk in chunks]
        
        # Store in ChromaDB
        collection.add(
            embeddings=embeddings,
            documents=chunk_texts,
            metadatas=metadatas,
            ids=ids
        )
        
        # Also store in Supabase for backup/search
        for i, chunk in enumerate(chunks):
            chunk_data = {
                "file_id": file_info['id'],
                "user_id": user_id,
                "chunk_index": chunk.chunk_index,
                "content": chunk.content,
                "embedding": embeddings[i],
                "metadata": chunk.metadata
            }
            
            get_supabase_client().table('document_chunks').insert(chunk_data).execute()
        
        logger.info(f"Stored {len(chunks)} chunks for file {file_info['name']}")
        
    except Exception as e:
        logger.error(f"Chunk storage error: {e}")

async def store_file_metadata(
    user_id: str,
    file_info: Dict,
    source: str,
    chunk_count: int
):
    """Store file metadata in database"""
    try:
        metadata = {
            "id": file_info['id'],
            "user_id": user_id,
            "name": file_info['name'],
            "path": file_info.get('webViewLink', ''),
            "size": int(file_info.get('size', 0)) if file_info.get('size') else 0,
            "mime_type": file_info.get('mimeType', ''),
            "source": source,
            "external_id": file_info['id'],
            "external_url": file_info.get('webViewLink', ''),
            "processed": True,
            "processing_status": "completed",
            "metadata": {
                "chunk_count": chunk_count,
                "modified_time": file_info.get('modifiedTime', file_info.get('lastModifiedDateTime', ''))
            }
        }
        
        # Upsert file metadata
        get_supabase_client().table('file_metadata').upsert(metadata).execute()
        
    except Exception as e:
        logger.error(f"File metadata storage error: {e}")

@router.get("/files/{user_id}")
async def get_processed_files(user_id: str):
    """Get list of processed files for user"""
    try:
        files_data = get_supabase_client().table('file_metadata').select('*').eq('user_id', user_id).eq('processed', True).execute()
        return {"files": files_data.data}
        
    except Exception as e:
        logger.error(f"Get files error: {e}")
        raise HTTPException(status_code=500, detail="Failed to get files")

@router.get("/status")
async def get_embed_status():
    """Get embedding status for compatibility with frontend"""
    try:
        # Return basic status - in a real app you'd check actual embedding jobs
        return {
            "status": "ready",
            "jobs_running": 0,
            "jobs_completed": 0,
            "vector_store_connected": True
        }
    except Exception as e:
        logger.error(f"Embed status error: {e}")
        raise HTTPException(status_code=500, detail="Failed to get embed status")

@router.get("/status/{user_id}")
async def get_user_embed_status(user_id: str):
    """Get embedding status for a specific user"""
    try:
        # Get job logs for this user
        jobs_data = get_supabase_client().table('job_logs').select('*').eq('user_id', user_id).order('created_at', desc=True).limit(10).execute()
        
        running_jobs = [job for job in jobs_data.data if job['status'] == 'running']
        completed_jobs = [job for job in jobs_data.data if job['status'] == 'completed']
        failed_jobs = [job for job in jobs_data.data if job['status'] == 'failed']
        
        return {
            "status": "running" if running_jobs else "ready",
            "jobs_running": len(running_jobs),
            "jobs_completed": len(completed_jobs),
            "jobs_failed": len(failed_jobs),
            "recent_jobs": jobs_data.data,
            "vector_store_connected": True
        }
        
    except Exception as e:
        logger.error(f"User embed status error: {e}")
        raise HTTPException(status_code=500, detail="Failed to get user embed status")