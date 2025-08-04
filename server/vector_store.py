"""
Briefly Solo - ChromaDB Vector Store Implementation
Replaces FAISS with ChromaDB for better Windows compatibility and easier installation
"""

import os
import json
import logging
import hashlib
from typing import List, Dict, Any, Optional
from pathlib import Path
import threading
import time
from dotenv import load_dotenv
# Optional imports for ML libraries (graceful degradation for serverless deployment)
try:
    import chromadb
    from chromadb.config import Settings
    from sentence_transformers import SentenceTransformer
    import numpy as np
    import torch
    ML_LIBRARIES_AVAILABLE = True
except ImportError as e:
    logger.warning(f"ML libraries not available: {e}")
    ML_LIBRARIES_AVAILABLE = False
    # Create dummy classes for compatibility
    class chromadb:
        pass
    class Settings:
        pass
    class SentenceTransformer:
        def __init__(self, *args, **kwargs):
            pass
    np = None
    torch = None

import traceback

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Global lock to prevent concurrent indexing
vector_index_lock = threading.Lock()

# File hash cache for incremental indexing
CACHE_FILE = "data/indexed_files.json"
class FileCache:
    def __init__(self, cache_file: str = CACHE_FILE):
        self.cache_file = cache_file
        self.data = self._load_cache()

    def _load_cache(self) -> Dict[str, Dict[str, Any]]:
        if os.path.exists(self.cache_file):
            try:
                with open(self.cache_file, 'r') as f:
                    return json.load(f)
            except Exception as e:
                logger.warning(f"Failed to load cache: {e}")
        return {}

    def save(self):
        try:
            with open(self.cache_file, 'w') as f:
                json.dump(self.data, f, indent=2)
        except Exception as e:
            logger.error(f"Failed to save cache: {e}")

    def get_hash(self, filepath: str) -> str:
        return self.data.get(filepath, {}).get("hash", "")

    def update_hash(self, filepath: str, file_hash: str, chunk_count: int):
        self.data[filepath] = {
            "hash": file_hash,
            "chunk_count": chunk_count
        }

    @staticmethod
    def compute_hash(filepath: str) -> str:
        h = hashlib.sha256()
        try:
            with open(filepath, "rb") as f:
                while chunk := f.read(8192):
                    h.update(chunk)
            return h.hexdigest()
        except Exception as e:
            logger.warning(f"Could not hash file {filepath}: {e}")
            return ""



class ChromaVectorStore:
    """
    ChromaDB-based vector store for document embeddings and similarity search.
    Provides the same interface as the previous FAISS implementation.
    """
    
    def __init__(self, persist_directory: str = "data/chroma_db", collection_name: str = "documents"):
        """
        Initialize ChromaDB vector store.
        
        Args:
            persist_directory: Directory to persist ChromaDB data
            collection_name: Name of the ChromaDB collection
        """
        if not ML_LIBRARIES_AVAILABLE:
            logger.warning("ChromaVectorStore initialized but ML libraries not available")
            self.available = False
            return
            
        self.available = True
        self.persist_directory = Path(persist_directory)
        self.persist_directory.mkdir(parents=True, exist_ok=True)
        self.collection_name = collection_name
        
        # Initialize ChromaDB client (Cloud or local)
        chroma_host = os.getenv("CHROMA_HOST")
        chroma_api_key = os.getenv("CHROMA_API_KEY")
        
        if chroma_host and chroma_api_key:
            # Use Chroma Cloud
            self.client = chromadb.HttpClient(
                host=chroma_host,
                port=int(os.getenv("CHROMA_PORT", 8000)),
                ssl=True,
                headers={
                    "X-Chroma-Token": chroma_api_key
                }
            )
            logger.info(f"[VECTOR] Using Chroma Cloud at {chroma_host}")
        else:
            # Fallback to local ChromaDB
            self.client = chromadb.PersistentClient(
                path=str(self.persist_directory),
                settings=Settings(
                    anonymized_telemetry=False,
                    allow_reset=True
                )
            )
            logger.info(f"[VECTOR] Using local ChromaDB at {self.persist_directory}")
        
        # Robust cross-platform device detection (CUDA for NVIDIA, MPS for Apple Silicon, fallback to CPU)
        torch_version = getattr(torch, '__version__', 'unknown')
        cuda_version = getattr(torch.version, 'cuda', 'unknown')
        cuda_available = torch.cuda.is_available()
        device_count = torch.cuda.device_count() if cuda_available else 0
        mps_available = hasattr(torch.backends, "mps") and torch.backends.mps.is_available() and torch.backends.mps.is_built()
        device = "cpu"
        device_name = "None"
        backend = "cpu"
        if cuda_available and device_count > 0:
            device = "cuda"
            device_name = torch.cuda.get_device_name(0)
            backend = f"CUDA {cuda_version}"
            logger.info(f"[VECTOR] Using NVIDIA GPU for embedding: {device_name} (CUDA {cuda_version})")
        elif mps_available:
            device = "mps"
            device_name = "Apple MPS"
            backend = "MPS"
            logger.info(f"[VECTOR] Using Apple Silicon GPU for embedding: {device_name}")
        else:
            logger.warning(f"[VECTOR] No GPU detected. Using CPU. If you have a GPU, install the correct PyTorch build. See https://pytorch.org/get-started/locally/")
        self.device = device
        self.device_name = device_name
        self.backend = backend
        self.model = SentenceTransformer('all-mpnet-base-v2', device=device)
        logger.info(f"[VECTOR] Using embedding model: all-mpnet-base-v2 on device: {device} ({device_name}, backend: {backend})")
        
        # Get or create collection
        try:
            self.collection = self.client.get_collection(name=self.collection_name)
            logger.info(f"Loaded existing collection '{self.collection_name}' with {self.collection.count()} documents")
        except Exception:
            self.collection = self.client.create_collection(
                name=self.collection_name,
                metadata={"description": "Briefly Solo document embeddings"}
            )
            logger.info(f"Created new collection '{self.collection_name}'")
    
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
        """
        Split a document into overlapping chunks for embedding.
        """
        import re
        tokens = re.findall(r"\S+", text)
        chunks = []
        start = 0
        while start < len(tokens):
            end = min(start + chunk_size, len(tokens))
            chunk = " ".join(tokens[start:end])
            chunks.append(chunk)
            if end == len(tokens):
                break
            start += chunk_size - overlap
        return chunks

    def add_documents(self, texts: List[str], metadatas: List[Dict[str, Any]] = None, ids: List[str] = None, chunk_size: int = 500, overlap: int = 100) -> None:
        """
        Add documents to the vector store, splitting into small embeddable chunks.
        Retries collection creation if missing.
        """
        if not texts:
            logger.warning("No texts provided to add_documents")
            return

        chunked_texts = []
        chunked_metadatas = []
        chunked_ids = []
        for i, text in enumerate(texts):
            meta = metadatas[i] if metadatas and i < len(metadatas) else {"source": f"document_{i}"}
            chunks = self.chunk_text(text, chunk_size=chunk_size, overlap=overlap)
            for j, chunk in enumerate(chunks):
                chunked_texts.append(chunk)
                chunk_meta = dict(meta)
                chunk_meta["chunk_index"] = j
                chunk_meta["total_chunks"] = len(chunks)
                chunked_metadatas.append(chunk_meta)
                chunked_ids.append(f"doc_{i}_chunk_{j}_{hash(chunk[:100])}")

        import time
        logger.info(f"Generating embeddings for {len(chunked_texts)} document chunks...")
        start_time = time.perf_counter()
        embeddings = self.model.encode(chunked_texts, show_progress_bar=True)
        elapsed = time.perf_counter() - start_time
        per_chunk = elapsed / len(chunked_texts) if chunked_texts else 0
        eta = per_chunk * len(chunked_texts)
        logger.info(f"[VECTOR] Indexed {len(chunked_texts)} chunks in {elapsed:.2f} seconds (~{per_chunk:.2f} sec/chunk)")
        logger.info(f"[VECTOR] Estimated total indexing time: {eta / 60:.1f} minutes")

        for attempt in range(2):
            try:
                self.collection.add(
                    embeddings=embeddings.tolist(),
                    documents=chunked_texts,
                    metadatas=chunked_metadatas,
                    ids=chunked_ids
                )
                logger.info(f"Successfully added {len(chunked_texts)} document chunks to vector store")
                return
            except Exception as e:
                logger.error(f"Error adding documents to vector store (attempt {attempt+1}): {e}\n{traceback.format_exc()}")
                if 'does not exist' in str(e) or 'not found' in str(e):
                    try:
                        self.collection = self.client.create_collection(
                            name=self.collection_name,
                            metadata={"description": "Briefly Solo document embeddings (auto-recreated)"}
                        )
                        logger.warning(f"Collection '{self.collection_name}' was missing and has been recreated.")
                        continue
                    except Exception as ce:
                        logger.error(f"Failed to recreate collection: {ce}")
                        raise
                else:
                    raise
        raise RuntimeError("Failed to add documents to vector store after retrying.")
    
    def similarity_search(self, query: str, k: int = 5, filter_dict: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """
        Search for similar documents.
        
        Args:
            query: Query text to search for
            k: Number of results to return
            filter_dict: Optional metadata filter
            
        Returns:
            List of dictionaries containing document text, metadata, and similarity scores
        """
        # Generate query embedding
        query_embedding = self.model.encode([query])
        
        # Perform search
        try:
            results = self.collection.query(
                query_embeddings=query_embedding.tolist(),
                n_results=k,
                where=filter_dict,
                include=["documents", "metadatas", "distances"]
            )
            
            # Format results
            formatted_results = []
            if results['documents'] and results['documents'][0]:
                for i, (doc, metadata, distance) in enumerate(zip(
                    results['documents'][0],
                    results['metadatas'][0],
                    results['distances'][0]
                )):
                    formatted_results.append({
                        'text': doc,
                        'metadata': metadata,
                        'similarity_score': 1 - distance,  # Convert distance to similarity
                        'rank': i + 1
                    })
            
            logger.info(f"Found {len(formatted_results)} similar documents for query")
            return formatted_results
            
        except Exception as e:
            logger.error(f"Error during similarity search: {e}")
            return []
    
    def get_collection_stats(self) -> Dict[str, Any]:
        """
        Get statistics about the vector store collection.
        
        Returns:
            Dictionary with collection statistics
        """
        try:
            count = self.collection.count()
            return {
                "collection_name": self.collection_name,
                "document_count": count,
                "doc_count": count,
                "persist_directory": str(self.persist_directory),
                "model_name": self.model.get_sentence_embedding_dimension(),
                "embedding_dimension": self.model.get_sentence_embedding_dimension()
            }
        except Exception as e:
            logger.error(f"Error getting collection stats: {e}")
            return {"error": str(e)}
    
    def delete_collection(self) -> bool:
        """
        Delete the entire collection.
        
        Returns:
            True if successful, False otherwise
        """
        try:
            self.client.delete_collection(name=self.collection_name)
            logger.info(f"Deleted collection '{self.collection_name}'")
            return True
        except Exception as e:
            logger.error(f"Error deleting collection: {e}")
            return False
    
    def reset_collection(self) -> bool:
        """
        Reset the collection (delete all documents but keep the collection).
        
        Returns:
            True if successful, False otherwise
        """
        try:
            # Delete and recreate collection
            self.client.delete_collection(name=self.collection_name)
            self.collection = self.client.create_collection(
                name=self.collection_name,
                metadata={"description": "Briefly Solo document embeddings"}
            )
            logger.info(f"Reset collection '{self.collection_name}'")
            return True
        except Exception as e:
            logger.error(f"Error resetting collection: {e}")
            return False
    
    def update_document(self, doc_id: str, text: str, metadata: Dict[str, Any] = None) -> bool:
        """
        Update a specific document in the collection.
        
        Args:
            doc_id: ID of the document to update
            text: New text content
            metadata: New metadata
            
        Returns:
            True if successful, False otherwise
        """
        try:
            # Generate new embedding
            embedding = self.model.encode([text])
            
            # Update in collection
            self.collection.update(
                ids=[doc_id],
                embeddings=embedding.tolist(),
                documents=[text],
                metadatas=[metadata] if metadata else None
            )
            logger.info(f"Updated document '{doc_id}'")
            return True
        except Exception as e:
            logger.error(f"Error updating document '{doc_id}': {e}\n{traceback.format_exc()}")
            return False
    
    def delete_documents(self, doc_ids: List[str]) -> bool:
        """
        Delete specific documents from the collection.
        
        Args:
            doc_ids: List of document IDs to delete
            
        Returns:
            True if successful, False otherwise
        """
        try:
            self.collection.delete(ids=doc_ids)
            logger.info(f"Deleted {len(doc_ids)} documents")
            return True
        except Exception as e:
            logger.error(f"Error deleting documents: {e}")
            return False

# Backward compatibility - maintain the same interface as FAISS implementation
class VectorStore(ChromaVectorStore):
    """
    Backward compatibility wrapper for the ChromaVectorStore.
    Maintains the same interface as the previous FAISS implementation.
    """
    
    def __init__(self, data_dir: str = "data"):
        """Initialize with backward compatible parameters."""
        persist_dir = os.path.join(data_dir, "chroma_db")
        super().__init__(persist_directory=persist_dir)
    
    def add_texts(self, texts: List[str], metadatas: List[Dict] = None):
        """Backward compatible method name."""
        return self.add_documents(texts, metadatas)
    
    def search(self, query: str, k: int = 5):
        """Backward compatible search method."""
        return self.similarity_search(query, k)

# Example usage and testing
if __name__ == "__main__":
    # Initialize vector store
    vector_store = ChromaVectorStore()
    
    # Test documents
    test_docs = [
        "Python is a high-level programming language.",
        "Machine learning is a subset of artificial intelligence.",
        "ChromaDB is a vector database for AI applications.",
        "FastAPI is a modern web framework for building APIs.",
        "Vector embeddings represent text as numerical vectors."
    ]
    
    test_metadata = [
        {"category": "programming", "language": "python"},
        {"category": "ai", "topic": "machine_learning"},
        {"category": "database", "type": "vector"},
        {"category": "web", "framework": "fastapi"},
        {"category": "ai", "topic": "embeddings"}
    ]
    
    # Add documents
    print("Adding test documents...")
    vector_store.add_documents(test_docs, test_metadata)
    
    # Test search
    print("\nTesting similarity search...")
    results = vector_store.similarity_search("What is machine learning?", k=3)
    
    for i, result in enumerate(results, 1):
        print(f"\nResult {i}:")
        print(f"Text: {result['text']}")
        print(f"Similarity: {result['similarity_score']:.3f}")
        print(f"Metadata: {result['metadata']}")
    
    # Get stats
    print(f"\nCollection stats: {vector_store.get_collection_stats()}")
def build_vector_index(folder_path: str, force_rebuild: bool = False) -> bool:
    """
    Build vector index from documents in a folder.
    Uses file hash cache to skip unchanged files for incremental indexing.
    Retains advanced parsing for all supported file types.
    
    Args:
        folder_path: Path to folder containing documents
        force_rebuild: Whether to rebuild even if index exists
        
    Returns:
        True if successful, False otherwise
    """
    # Check if ML libraries are available
    if not ML_LIBRARIES_AVAILABLE:
        logger.warning("ML libraries not available, vector indexing disabled")
        return False
        
    try:
        import os
        from pathlib import Path
        
        logger.info(f"[VECTOR INDEX] Starting indexing for folder: {folder_path} (force_rebuild={force_rebuild})")
        # Initialize vector store and cache
        vector_store = ChromaVectorStore()
        file_cache = FileCache()
        
        # If force rebuild, reset the collection and clear cache
        if force_rebuild:
            vector_store.reset_collection()
            file_cache.data = {}
            file_cache.save()

        folder = Path(folder_path)
        if not folder.exists() or not folder.is_dir():
            logger.error(f"Folder does not exist: {folder_path}")
            return False

        # Supported file extensions for each type
        text_extensions = {'.txt', '.md', '.py', '.js', '.ts', '.json', '.csv', '.html', '.css'}
        pdf_extensions = {'.pdf'}
        docx_extensions = {'.docx'}
        xlsx_extensions = {'.xlsx'}
        xls_extensions = {'.xls'}
        pptx_extensions = {'.pptx'}
        ppt_extensions = {'.ppt'}
        rtf_extensions = {'.rtf'}
        odt_extensions = {'.odt'}
        ods_extensions = {'.ods'}
        odp_extensions = {'.odp'}
        epub_extensions = {'.epub'}
        xml_extensions = {'.xml'}
        doc_extensions = {'.doc'}  # Legacy Word, best effort

        texts = []
        metadatas = []

        for file_path in folder.rglob('*'):
            if not file_path.is_file():
                continue
            ext = file_path.suffix.lower()
            content = None
            # --- File hash cache logic ---
            file_hash = FileCache.compute_hash(str(file_path))
            if not force_rebuild and file_cache.get_hash(str(file_path)) == file_hash:
                logger.info(f"[CACHE] Skipping unchanged file: {file_path}")
                continue
            try:
                if ext in text_extensions:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                elif ext in pdf_extensions:
                    import pdfplumber
                    with pdfplumber.open(str(file_path)) as pdf:
                        content = "\n".join(page.extract_text() or '' for page in pdf.pages)

                elif ext in docx_extensions:
                    from docx import Document
                    doc = Document(str(file_path))
                    content = "\n".join([p.text for p in doc.paragraphs])

                elif ext in xlsx_extensions:
                    import openpyxl
                    try:
                        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
                        content = ""
                        for sheet in wb.sheetnames:
                            ws = wb[sheet]
                            rows = list(ws.iter_rows(values_only=True))
                            if not rows:
                                continue
                            # Convert to markdown table
                            header = rows[0]
                            table = ["| " + " | ".join(str(h) if h is not None else "" for h in header) + " |",
                                     "| " + " | ".join(["---"] * len(header)) + " |"]
                            for row in rows[1:]:
                                table.append("| " + " | ".join(str(cell) if cell is not None else "" for cell in row) + " |")
                            content += f"\nSheet: {sheet}\n" + "\n".join(table) + "\n"
                        logger.info(f"Parsed Excel file {file_path} with {len(wb.sheetnames)} sheets.")
                    except Exception as e:
                        logger.warning(f"Failed to parse Excel file {file_path}: {e}")
                        content = ""

                elif ext in xls_extensions:
                    import xlrd
                    wb = xlrd.open_workbook(str(file_path))
                    content = "\n".join(
                        "\n".join(str(cell.value) for cell in sheet.row(row_idx) if cell.value)
                        for sheet in wb.sheets() for row_idx in range(sheet.nrows)
                    )

                elif ext in pptx_extensions:
                    from pptx import Presentation
                    prs = Presentation(str(file_path))
                    slides = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slides.append(shape.text)
                    content = "\n".join(slides)

                elif ext in ppt_extensions:
                    # Legacy PPT not supported natively, skip or log
                    logger.warning(f"Legacy .ppt file not supported for {file_path}")
                    continue

                elif ext in rtf_extensions:
                    from striprtf.striprtf import rtf_to_text
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = rtf_to_text(f.read())

                elif ext in odt_extensions or ext in ods_extensions or ext in odp_extensions:
                    from odf.opendocument import load
                    from odf import text as odf_text
                    doc = load(str(file_path))
                    paragraphs = doc.getElementsByType(odf_text.P)
                    content = "\n".join([p.firstChild.data if p.firstChild else '' for p in paragraphs])

                elif ext in epub_extensions:
                    from ebooklib import epub
                    book = epub.read_epub(str(file_path))
                    content = "\n".join(
                        item.get_content().decode('utf-8', errors='ignore')
                        for item in book.get_items() if item.get_type() == epub.ITEM_DOCUMENT
                    )

                elif ext in xml_extensions:
                    import xml.etree.ElementTree as ET
                    tree = ET.parse(str(file_path))
                    root = tree.getroot()
                    content = ET.tostring(root, encoding='unicode', method='text')

                elif ext in doc_extensions:
                    # Legacy .doc not supported natively, skip or log
                    logger.warning(f"Legacy .doc file not supported for {file_path}")
                    continue

                else:
                    # Unknown file type, skip
                    logger.info(f"Skipping unsupported file type: {file_path}")
                    continue

                if content and content.strip():
                    texts.append(content)
                    metadatas.append({
                        'filename': file_path.name,
                        'filepath': str(file_path.relative_to(folder)),
                        'extension': ext,
                        'size': len(content)
                    })
            except Exception as e:
                logger.warning(f"Could not read file {file_path}: {e}")
                continue

        if not texts:
            logger.warning(f"No readable files found in {folder_path}")
            return False

        # Acquire global lock to prevent concurrent indexing
        logger.info("Waiting for vector index lock...")
        acquired = vector_index_lock.acquire(timeout=600)  # 10 min timeout
        if not acquired:
            logger.error("Timeout waiting for vector index lock. Another indexing operation is running.")
            return False
        logger.info("Vector index lock acquired.")
        try:
            # Add documents to vector store (with chunking)
            logger.info(f"Building vector index from {len(texts)} documents (with chunking)...")
            vector_store.add_documents(texts, metadatas)

            logger.info(f"Successfully built vector index with {len(texts)} source documents (chunked)")
            return True
        finally:
            vector_index_lock.release()
            logger.info("Vector index lock released.")
        
    except Exception as e:
        logger.error(f"Error during indexing: {e}\n{traceback.format_exc()}")
        return False

from spellchecker import SpellChecker


def get_relevant_context(query: str, max_results: int = 1, top_k: int | None = None) -> List[Dict[str, Any]]:
    """
    Get relevant context for a query from the vector store.
    
    Args:
        query: Query text to search for
        max_results: Maximum number of results to return
        
    Returns:
        List of relevant documents with metadata
    """
    try:
        # Initialize vector store
        vector_store = ChromaVectorStore()

        # ---------------------------
        # 1) Autocorrect the query
        # ---------------------------
        spell = SpellChecker()
        corrected_query = " ".join(
            spell.correction(w) or w for w in query.split()
        )
        # Use corrected query for retrieval, but fall back to original if identical
        retrieval_query = corrected_query if corrected_query.lower() != query.lower() else query

        # Determine how many results to retrieve before re-ranking
        top_k = top_k or max(max_results * 4, 10)

        # Perform similarity search
        results = vector_store.similarity_search(retrieval_query, k=top_k)

        # If the autocorrected query produced no hits, try again with the original (just in case)
        if not results:
            results = vector_store.similarity_search(query, k=top_k)

        # ---------------------------
        # 2) Re-rank with keyword presence
        # ---------------------------
        corrected_tokens = [w.lower() for w in corrected_query.split()]
        for r in results:
            text_lower = r['text'].lower()
            keyword_hits = sum(1 for t in corrected_tokens if t in text_lower)
            r['keyword_hits'] = keyword_hits
        # Sort: similarity first, then keyword hits
        results.sort(key=lambda r: (r['similarity_score'], r['keyword_hits']), reverse=True)

        # Trim to desired max_results
        results = results[:max_results]
        
        # Format results for backward compatibility
        formatted_results = []
        for result in results:
            formatted_results.append({
                'content': result['text'],
                'metadata': result['metadata'],
                'score': result['similarity_score'],
                'rank': result['rank']
            })
        
        logger.info(f"[VECTOR] Retrieved {len(formatted_results)} relevant documents for query: '{query}'")
        for i, chunk in enumerate(formatted_results):
            logger.debug(
                f"[VECTOR] Chunk {i+1}:\n"
                f"  Score: {chunk['score']:.4f}\n"
                f"  Metadata: {chunk['metadata']}\n"
                f"  Preview: {chunk['content'][:200]}{'...' if len(chunk['content']) > 200 else ''}"
            )
        return formatted_results
        
    except Exception as e:
        logger.error(f"[VECTOR] Error getting relevant context: {e}\n{traceback.format_exc()}")
        return []

def get_vector_store_stats() -> Dict[str, Any]:
    """
    Get statistics about the current vector store.
    
    Returns:
        Dictionary with vector store statistics
    """
    # Check if ML libraries are available
    if not ML_LIBRARIES_AVAILABLE:
        return {
            "status": "disabled",
            "reason": "ML libraries not available",
            "collections": 0,
            "total_documents": 0
        }
        
    try:
        vector_store = ChromaVectorStore()
        return vector_store.get_collection_stats()
    except Exception as e:
        logger.error(f"Error getting vector store stats: {e}\n{traceback.format_exc()}")
        return {"error": str(e)}

